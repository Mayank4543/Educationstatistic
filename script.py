from pptx import Presentation
from pptx.util import Inches

# Create a presentation object
prs = Presentation()

# Slide 1: Title Slide
slide_layout = prs.slide_layouts[0]
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
subtitle = slide.placeholders[1]

title.text = "Educational Statistics - Data Analysis & Prediction"
subtitle.text = "A Data-Driven Approach to Understanding Educational Metrics"

# Slide 2: Description of the Tool
slide_layout = prs.slide_layouts[1]
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
content = slide.placeholders[1]

title.text = "Description of the Tool"
content.text = ("This tool analyzes educational data using visualization and prediction models. "
                "It helps understand trends in literacy rates, pupil-teacher ratios, and school availability. "
                "Using machine learning models, it predicts key metrics to assist in policy-making.")

# Slide 3: Algorithms/Tools Used for Prediction & Visualization
slide_layout = prs.slide_layouts[1]
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
content = slide.placeholders[1]

title.text = "Algorithms & Tools Used"
content.text = ("🔹 **Prediction Models:** Linear Regression, Random Forest, Decision Tree\n"
                "🔹 **Visualization:** Matplotlib, Seaborn, Plotly\n"
                "🔹 **Libraries Used:** Pandas, NumPy, Scikit-learn\n"
                "🔹 **Data Source:** Data.gov.in - Education Statistics Dataset")

# Slide 4: Outputs
slide_layout = prs.slide_layouts[1]
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
content = slide.placeholders[1]

title.text = "Model Outputs & Visualizations"
content.text = ("📌 **Heatmaps & Correlation Analysis** - Understanding feature relationships\n"
                "📌 **Regression Model Accuracy** - Mean Absolute Error (MAE) & R² Score\n"
                "📌 **Predictions on Educational Metrics** - Literacy Rate, Enrollment Ratios\n"
                "📌 **Feature Importance Analysis** - Identifying key influencing factors")

# Slide 5: Conclusion
slide_layout = prs.slide_layouts[1]
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
content = slide.placeholders[1]

title.text = "Conclusion"
content.text = ("✔️ Data-driven insights help in understanding educational gaps.\n"
                "✔️ Prediction models provide future estimates for informed decision-making.\n"
                "✔️ Visualizations offer intuitive understanding for stakeholders.\n"
                "✔️ Further improvements include better feature selection & advanced models.")

# Save the presentation
pptx_file = "/mnt/data/Educational_Statistics_Prediction.pptx"
prs.save(pptx_file)

pptx_file
from pptx import Presentation
from pptx.util import Inches

# Create a presentation object
prs = Presentation()

# Slide 1: Title Slide
slide_layout = prs.slide_layouts[0]
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
subtitle = slide.placeholders[1]

title.text = "Educational Statistics - Data Analysis & Prediction"
subtitle.text = "A Data-Driven Approach to Understanding Educational Metrics"

# Slide 2: Description of the Tool
slide_layout = prs.slide_layouts[1]
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
content = slide.placeholders[1]

title.text = "Description of the Tool"
content.text = ("This tool analyzes educational data using visualization and prediction models. "
                "It helps understand trends in literacy rates, pupil-teacher ratios, and school availability. "
                "Using machine learning models, it predicts key metrics to assist in policy-making.")

# Slide 3: Algorithms/Tools Used for Prediction & Visualization
slide_layout = prs.slide_layouts[1]
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
content = slide.placeholders[1]

title.text = "Algorithms & Tools Used"
content.text = ("🔹 **Prediction Models:** Linear Regression, Random Forest, Decision Tree\n"
                "🔹 **Visualization:** Matplotlib, Seaborn, Plotly\n"
                "🔹 **Libraries Used:** Pandas, NumPy, Scikit-learn\n"
                "🔹 **Data Source:** Data.gov.in - Education Statistics Dataset")

# Slide 4: Outputs
slide_layout = prs.slide_layouts[1]
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
content = slide.placeholders[1]

title.text = "Model Outputs & Visualizations"
content.text = ("📌 **Heatmaps & Correlation Analysis** - Understanding feature relationships\n"
                "📌 **Regression Model Accuracy** - Mean Absolute Error (MAE) & R² Score\n"
                "📌 **Predictions on Educational Metrics** - Literacy Rate, Enrollment Ratios\n"
                "📌 **Feature Importance Analysis** - Identifying key influencing factors")

# Slide 5: Conclusion
slide_layout = prs.slide_layouts[1]
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
content = slide.placeholders[1]

title.text = "Conclusion"
content.text = ("✔️ Data-driven insights help in understanding educational gaps.\n"
                "✔️ Prediction models provide future estimates for informed decision-making.\n"
                "✔️ Visualizations offer intuitive understanding for stakeholders.\n"
                "✔️ Further improvements include better feature selection & advanced models.")

# Save the presentation
pptx_file = "C:/Users/Mayank Rathore/Desktop/Educational_Statistics_Prediction.pptx"


prs.save(pptx_file)

pptx_file
